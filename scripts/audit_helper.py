#!/usr/bin/env python3
"""
audit_helper.py — mechanical pre-checks for the Deck Audit mode.

This does NOT replace the judgment-based scoring in references/audit-rubric.md.
It catches the checkable, objective stuff cheaply first, so the model can spend
its judgment on structure and messaging instead of re-deriving "slide 7 has too
much text" by eye.

Usage:
    python3 audit_helper.py deck.pptx

Requires: python-pptx (pip install python-pptx --break-system-packages)
"""

import sys
import re
from collections import Counter

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is required: pip install python-pptx --break-system-packages")
    sys.exit(1)

# Rough density guideline: past this many characters of body text, a slide is
# very likely being read rather than glanced at. This is a heuristic trigger
# for review, not a hard rule — dense reference slides can be legitimate.
DENSE_SLIDE_CHAR_THRESHOLD = 150

# Topic-label titles tend to be short noun phrases; action titles tend to be
# longer sentences that state a takeaway. This is a soft signal, not a verdict —
# always read the actual title before concluding anything.
SHORT_TITLE_CHAR_THRESHOLD = 8

PLACEHOLDER_PATTERNS = re.compile(
    r"lorem ipsum|\btodo\b|\bxxx+\b|\btbd\b|\[insert|请替换|待补充|占位|示例文本",
    re.IGNORECASE,
)


def char_count(text: str) -> int:
    # Count Chinese/Japanese/Korean characters as visually "full" characters;
    # everything else (Latin, digits, punctuation) counts per-character too —
    # this is a density proxy, not a precise typographic measure.
    return len(text)


def analyze(path: str):
    prs = Presentation(path)
    dense_slides = []
    short_titles = []
    placeholder_hits = []
    fonts_used = Counter()
    slide_reports = []

    for i, slide in enumerate(prs.slides, start=1):
        all_text = []
        title_text = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            frame_text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
            if not frame_text:
                continue
            all_text.append(frame_text)

            is_title_placeholder = getattr(shape, "is_placeholder", False) and getattr(
                shape.placeholder_format, "type", None
            ) is not None and "TITLE" in str(shape.placeholder_format.type)
            if is_title_placeholder and title_text is None:
                title_text = frame_text.strip()

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font and run.font.name:
                        fonts_used[run.font.name] += 1

        full_text = "\n".join(all_text)
        n_chars = char_count(full_text)

        if n_chars > DENSE_SLIDE_CHAR_THRESHOLD:
            dense_slides.append((i, n_chars))

        if title_text and len(title_text) <= SHORT_TITLE_CHAR_THRESHOLD:
            short_titles.append((i, title_text))

        if PLACEHOLDER_PATTERNS.search(full_text):
            placeholder_hits.append(i)

        slide_reports.append((i, title_text or "(no detected title)", n_chars))

    print(f"=== Deck Audit — mechanical pre-check: {path} ===\n")
    print(f"Total slides: {len(slide_reports)}\n")

    print("-- Per-slide overview --")
    for i, title, n_chars in slide_reports:
        flag = " ⚠ dense" if n_chars > DENSE_SLIDE_CHAR_THRESHOLD else ""
        print(f"  [{i:>2}] {n_chars:>4} chars  {title}{flag}")

    print("\n-- Font usage --")
    if fonts_used:
        for font, count in fonts_used.most_common():
            print(f"  {font}: {count} runs")
        if len(fonts_used) > 2:
            print(f"  ⚠ {len(fonts_used)} distinct fonts detected — check against chinese-typography.md (max 2 recommended)")
    else:
        print("  No explicit font names set on runs (inheriting from theme/layout — check visually).")

    print("\n-- Possible short/topic-label titles (verify manually) --")
    if short_titles:
        for i, t in short_titles:
            print(f"  [{i}] \"{t}\" — check whether this should be an action title instead")
    else:
        print("  None flagged.")

    print("\n-- Possible leftover placeholder text --")
    if placeholder_hits:
        for i in placeholder_hits:
            print(f"  [{i}] contains placeholder-like text (lorem/TODO/待补充/占位/etc.) — verify and fix")
    else:
        print("  None flagged.")

    print("\n-- Dense slides (over threshold, may need trimming or splitting) --")
    if dense_slides:
        for i, n in dense_slides:
            print(f"  [{i}] {n} characters of text")
    else:
        print("  None flagged.")

    print(
        "\nThese are mechanical signals only. Score the deck against "
        "references/audit-rubric.md using these as a starting point, not a substitute "
        "for reading it."
    )


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python3 audit_helper.py deck.pptx")
        sys.exit(1)
    analyze(sys.argv[1])
